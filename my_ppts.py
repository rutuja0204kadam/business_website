from wand.image import Image
from wand.compat import nested
import os
from pptx.util import Inches
from pptx import Presentation
from PIL import Image as pillow_image

# getlist of all the images in the directory
file_list = os.listdir(r"/home/rutujakadam/Downloads/my_ppt_folder")


# creating watermark image using ImageMagick
def watermark_image(the_image, pk, logo):

    # code to put watermark of logo on image
    image = Image(filename="/home/rutujakadam/Downloads/my_ppt_folder/"+the_image)
    image.watermark(logo, 0.0, 40, 40)

    # give a new name to watermarked image
    output_watermark_image_name = "out"+str(pk)+".png"

    # save the watermarked image in same directory as old image or you can even keep this watermarked image on old image
    # by saying output_watermark_image_name = "image"+str(pk+1)+".png"
    image.save(filename='/home/rutujakadam/Downloads/my_ppt_folder/' + output_watermark_image_name)
    return output_watermark_image_name


# A function that outputs the final ppt file
def create_presentation():

    # starting presentation
    prs = Presentation()

    # Resizing the logo image by keeping same aspect ratio using ImageMagick
    logo = Image(filename='/home/rutujakadam/Downloads/my_ppt_folder/nike_black.png')
    new_width = logo.width * 0.5
    logo.resize(int(new_width), int((logo.height/logo.width)*new_width))
    logo.save(filename='/home/rutujakadam/Downloads/my_ppt_folder/nike_black.png')


    # initializing for loop for passing all the images in directory
    for i in range(0, len(file_list)-1):

        # creating the slides in ppt
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # adding title
        title_shape = slide.shapes.title
        title_shape.text = 'Slide Title' + str(i + 1)

        # adding subtitle
        body_shape = slide.shapes.placeholders[1]
        sub_title = body_shape.text_frame
        sub_title.text = 'Slide Subtitle'+str(i+1)

        # going to function that watermarks the logo on images
        a = watermark_image(file_list[i], i, logo)

        # finding the aspect ratio of the Images with watermark
        image_path = '/home/rutujakadam/Downloads/my_ppt_folder/'+a
        img = pillow_image.open(image_path)
        w, h = img.size
        ratio = h / w

        # Adding the Images with watermark to the presentation keeping its aspect ratio
        # same and converting height and width from px to inches
        pic = slide.shapes.add_picture(image_path, left=Inches(1), top=Inches(2.5), height=Inches(ratio*2.8), width=Inches(2.8))

    # Inserting the ppt into same directory
    prs.save('/home/rutujakadam/Downloads/my_ppt_folder/test2.pptx')


# running the function to create the presentation
create_presentation()

